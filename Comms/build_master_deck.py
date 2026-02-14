# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette — same brand palette as EQT deck
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUBTLE_BG = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
TABLE_ROW_ALT = RGBColor(0xF0, 0xF4, 0xF8)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT = RGBColor(0xCC, 0x33, 0x33)
GREEN_ACCENT = RGBColor(0x33, 0xAA, 0x55)
FONT = "Calibri"


def add_dark_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()


def add_light_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()


def add_accent_bar(slide, y=Inches(0), height=Inches(0.06)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()


def set_cell_text(cell, text, font_size=11, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_footnote(slide, text, y=Inches(7.0)):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(8)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT
    run.font.italic = True


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_title(slide, text, color=DARK_BG):
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.color.rgb = color
    run.font.name = FONT
    run.font.bold = True


def add_subtitle(slide, text, y=Inches(1.4)):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.color.rgb = DARK_TEXT
    run.font.name = FONT


def add_callout(slide, text, y=Inches(6.0), accent_color=ACCENT_TEAL):
    """Add a callout box with left accent border."""
    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(0.06), Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Text
    txBox = slide.shapes.add_textbox(Inches(1.1), y, Inches(11.0), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(15)
    run.font.color.rgb = DARK_TEXT
    run.font.name = FONT
    run.font.italic = True


# ============================================================
# SLIDE 1 — Title [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, y=Inches(3.4), height=Inches(0.04))

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "realfast"
run.font.size = Pt(48)
run.font.color.rgb = WHITE
run.font.name = FONT
run.font.bold = True

txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10), Inches(0.6))
tf2 = txBox2.text_frame
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "AI-native IT services"
run.font.size = Pt(22)
run.font.color.rgb = ACCENT_TEAL
run.font.name = FONT

txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.5))
tf3 = txBox3.text_frame
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "[Presenter Name]  |  [Date]"
run.font.size = Pt(14)
run.font.color.rgb = MID_GRAY
run.font.name = FONT

add_speaker_notes(slide, "None. Let the slide sit for 3-5 seconds.")


# ============================================================
# SLIDE 2 — Governing Question [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "What does an IT services firm need to look like\nin three years to survive?"
run.font.size = Pt(36)
run.font.color.rgb = WHITE
run.font.name = FONT
run.font.bold = True

txBox2 = slide.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(11.0), Inches(0.6))
tf2 = txBox2.text_frame
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "$1T+ industry. Flat unit economics. Structural repricing underway."
run.font.size = Pt(18)
run.font.color.rgb = LIGHT_GRAY
run.font.name = FONT

add_speaker_notes(slide, "This is the question we've built our company around. Not 'how do we use AI' \u2014 that's a tool question. The structural question is: why has a trillion-dollar industry not improved its fundamental economics in ten years, and what happens when AI forces the issue? Everything I'm about to show you flows from this.")


# ============================================================
# SLIDE 3 — RPE Data [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "IT services revenue per employee has been flat to declining for a decade")

rpe_data = [
    ["Company", "RPE (FY2015)", "RPE (FY2023)", "Change"],
    ["Infosys", "$49,442", "$53,060", "+7.3% over 8 years"],
    ["TCS", "~$50,700", "~$47,000", "Flat to negative"],
    ["Wipro", "~$48,000", "~$45,100", "Declining"],
    ["HCLTech", "~$65,000", "~$61,400", "Declining"],
]

table = slide.shapes.add_table(len(rpe_data), 4, Inches(0.8), Inches(1.8), Inches(11.0), Inches(2.2)).table
col_widths = [Inches(2.5), Inches(2.8), Inches(2.8), Inches(2.9)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

for row_idx, row in enumerate(rpe_data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        is_header = row_idx == 0
        is_change = col_idx == 3 and row_idx > 0
        color = WHITE if is_header else (RED_ACCENT if is_change and "Declining" in val or is_change and "negative" in val else DARK_TEXT)
        set_cell_text(cell, val, font_size=13, bold=is_header or col_idx == 0, color=color)
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_WHITE

# TCS callout
add_callout(slide, "TCS: Revenue doubled ($16B \u2192 $30B). Headcount doubled (319K \u2192 608K). Ratio unchanged.", y=Inches(4.4))

add_footnote(slide, "Sources: Bullfincher, StockAnalysis, Pathfinders Training, Communications Today")

add_speaker_notes(slide, "TCS is the clearest case. They nearly doubled revenue over a decade. Sounds impressive \u2014 until you see headcount nearly doubled too. At roughly $50K RPE, every additional million dollars of revenue requires 20 new people. This is arithmetic, not opinion. When RPE does tick up \u2014 as it did for Infosys recently \u2014 it's driven by layoffs, not structural improvement. The top five firms shed 57,900 employees in two years. They'll rehire when deals return. RPE will revert.")


# ============================================================
# SLIDE 4 — Bench Eliminated [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "The bench that absorbed delivery risk has been systematically eliminated")

metrics = [
    ("Bench", "10-15%", "2-5%", "\u2193"),
    ("Utilization", "70-75%", "83-86%", "\u2191"),
    ("TCS Rule", "", "35 days", "Retrain or exit"),
]

card_w = Inches(3.4)
card_h = Inches(2.5)
start_x = Inches(0.8)

for i, (label, before, after, indicator) in enumerate(metrics):
    x = start_x + i * (card_w + Inches(0.25))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.3)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(20)
    run.font.color.rgb = DARK_BG
    run.font.name = FONT
    run.font.bold = True

    if before:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(14)
        run2 = p2.add_run()
        run2.text = f"{before}  \u2192  {after}"
        run2.font.size = Pt(24)
        run2.font.color.rgb = RED_ACCENT
        run2.font.name = FONT
        run2.font.bold = True
    else:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(14)
        run2 = p2.add_run()
        run2.text = after
        run2.font.size = Pt(28)
        run2.font.color.rgb = RED_ACCENT
        run2.font.name = FONT
        run2.font.bold = True

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(8)
    run3 = p3.add_run()
    run3.text = indicator if label == "TCS Rule" else ""
    run3.font.size = Pt(13)
    run3.font.color.rgb = MID_GRAY
    run3.font.name = FONT

add_footnote(slide, "Sources: Trak.in, Economic Times, NiftyTrader")

add_speaker_notes(slide, "Services firms used to absorb delivery risk through bench \u2014 idle people you could throw at a project when it slipped. That buffer is gone. Every project now has to hit its time-to-value target with no margin for error and no spare people. This isn't a choice they made for efficiency. It's a response to margin pressure. The safety net has been cut to protect the economics \u2014 but the economics haven't actually improved.")


# ============================================================
# SLIDE 5 — Deal Compression [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Customers want $150K AI projects \u2014 the firms serving them can't deliver below $3M profitably")

forces = [
    ("Shorter commitments", "Annual renewals replacing multi-year lock-ins"),
    ("Smaller, faster AI engagements", "$150K-$300K alongside multi-million dollar contracts"),
    ("Overhead per deal is constant", "The cost of processing a $300K deal approaches the cost of processing a $30M deal"),
]

card_w = Inches(3.5)
card_h = Inches(1.6)
start_x = Inches(0.8)

for i, (title, desc) in enumerate(forces):
    x = start_x + i * (card_w + Inches(0.2))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{i+1}. {title}"
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_BG
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(13)
    run2.font.color.rgb = MID_GRAY
    run2.font.name = FONT

add_callout(slide, "Firms must serve deals from $30K to $300M. Their cost structure makes them unprofitable below a threshold most don't know.", y=Inches(4.2), accent_color=RED_ACCENT)

add_speaker_notes(slide, "A sales engineer at a major cloud platform told us this directly: their largest enterprise client now regularly requests $150K-$300K AI projects alongside existing multi-million dollar contracts. The firm cannot profitably serve them. But the client implies a bigger deal will follow, so refusing isn't an option. They're caught between a cost structure they can't flex and a market they can't ignore. This is not one firm's problem. This is the industry.")


# ============================================================
# SLIDE 6 — Market Repricing [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "This is not a prediction \u2014 the public markets are pricing it in now")

# Left side — stock drops
stocks = [
    ("Cognizant", "~-10%"),
    ("Infosys ADRs", "-8.4%"),
    ("Nifty IT", "-6%"),
    ("Wipro", "-4.5%"),
]

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Stock reactions on AI agent launches (single session)"
run.font.size = Pt(14)
run.font.color.rgb = MID_GRAY
run.font.name = FONT
run.font.bold = True

for i, (name, drop) in enumerate(stocks):
    y = Inches(2.2) + i * Inches(0.9)
    # Bar background
    bar_width = Inches(abs(float(drop.replace("~", "").replace("-", "").replace("%", ""))) * 0.45)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, bar_width, Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED_ACCENT
    bar.line.fill.background()
    # Label
    txBox = slide.shapes.add_textbox(Inches(0.9), y + Inches(0.08), bar_width - Inches(0.1), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{name}: {drop}"
    run.font.size = Pt(14)
    run.font.color.rgb = WHITE
    run.font.name = FONT
    run.font.bold = True

# Right side — hiring collapse
right_data = [
    "TCS cut 11,000 in one quarter",
    "Big Four historically hired 10,000+/quarter",
    "Revenue growth: 15-20% \u2192 low single digits",
]

txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Hiring collapse"
run.font.size = Pt(14)
run.font.color.rgb = MID_GRAY
run.font.name = FONT
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.5), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(right_data):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    run = p.add_run()
    run.text = line
    run.font.size = Pt(18)
    run.font.color.rgb = DARK_TEXT
    run.font.name = FONT
    p.space_before = Pt(16)

add_speaker_notes(slide, "When capable AI agents launched in early 2025, the market's response was immediate. The entire IT outsourcing model is labour arbitrage \u2014 hire at $15/hour, bill at $80/hour. The market saw AI agents and immediately repriced the stocks. Revenue won't fall off a cliff \u2014 multi-year contracts protect the near term. But the trajectory is clear: eight to ten quarters of consensus estimates trimmed by two to three percent each cycle. A slow bleed, not a sudden collapse. The firms know this. Their stock prices know this. Their hiring plans know this.")


# ============================================================
# SLIDE 7 — Bottleneck Relocation [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Making engineers faster doesn't make the firm faster \u2014 the bottleneck just moves")

# Pipeline stages
stages = ["Pre-sales", "Legal", "Compliance", "Staffing", "ENGINEERING", "QA", "PM", "Delivery"]
stage_w = Inches(1.3)
stage_h = Inches(0.7)
start_x = Inches(0.5)
pipeline_y = Inches(2.5)

for i, stage in enumerate(stages):
    x = start_x + i * (stage_w + Inches(0.15))
    is_eng = stage == "ENGINEERING"
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, pipeline_y, stage_w, stage_h)
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT_BLUE if is_eng else RGBColor(0xE0, 0xE0, 0xE0)
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = stage
    run.font.size = Pt(11)
    run.font.color.rgb = WHITE if is_eng else DARK_TEXT
    run.font.name = FONT
    run.font.bold = is_eng

    # Arrow between stages
    if i < len(stages) - 1:
        arrow_x = x + stage_w
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, pipeline_y + Inches(0.2), Inches(0.15), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

# Callout below pipeline
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Grant 10x engineering speed. The bottleneck moves to legal.\nThen compliance. Then pre-sales.\nFriction per unit of revenue stays the same."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = FONT

add_speaker_notes(slide, "The standard response from every IT services firm: give engineers AI tools. Copilot, Cursor, Devin. Announce AI transformation. But this addresses one node in a multi-node pipeline. Even granting the most optimistic productivity claims \u2014 2x, 5x, 10x \u2014 everything upstream and downstream stays the same speed. The firm doesn't get faster. Individual engineers get faster. The distinction matters. The question isn't whether AI tools help developers. Assume they do. The question is whether developer productivity was ever the binding constraint. In most IT services engagements, it was not.")


# ============================================================
# SLIDE 8 — Friction Per Unit of Revenue [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "35-50% of delivery cost is organizational friction, not engineering")

# Engineering bar
eng_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.2))
eng_bar.fill.solid()
eng_bar.fill.fore_color.rgb = ACCENT_BLUE
eng_bar.line.fill.background()
tf = eng_bar.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Engineering: 50-65%"
run.font.size = Pt(22)
run.font.color.rgb = WHITE
run.font.name = FONT
run.font.bold = True

# Friction components
friction_items = ["Selling & scoping", "Legal review", "Compliance", "Staffing", "Project management", "Knowledge transfer", "Documentation & renewal"]
friction_w = Inches(1.55)
friction_h = Inches(0.5)

for i, item in enumerate(friction_items):
    col = i % 4
    row = i // 4
    x = Inches(6.8) + col * (friction_w + Inches(0.1))
    y = Inches(2.0) + row * (friction_h + Inches(0.1))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, friction_w, friction_h)
    box.fill.solid()
    box.fill.fore_color.rgb = SUBTLE_BG
    box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = item
    run.font.size = Pt(10)
    run.font.color.rgb = DARK_TEXT
    run.font.name = FONT

# Friction label
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(3.2), Inches(5.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Organizational friction: 35-50%"
run.font.size = Pt(18)
run.font.color.rgb = RED_ACCENT
run.font.name = FONT
run.font.bold = True

# Definition box
defn_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.9))
defn_box.fill.solid()
defn_box.fill.fore_color.rgb = DARK_BG
defn_box.line.fill.background()
tf = defn_box.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Friction per unit of revenue"
run.font.size = Pt(16)
run.font.color.rgb = ACCENT_TEAL
run.font.name = FONT
run.font.bold = True
run2 = p.add_run()
run2.text = " = total cost (time + money) to move $1 from first conversation to delivered outcome"
run2.font.size = Pt(16)
run2.font.color.rgb = LIGHT_GRAY
run2.font.name = FONT

add_footnote(slide, "Source: APQC benchmarks; Accenture SG&A ~17% of revenue before delivery costs")

add_speaker_notes(slide, "Revenue per employee tells you how productive your people are. It doesn't tell you how efficient your organization is at converting opportunity into delivered revenue. The metric that matters is friction per unit of revenue. For a typical IT services project, engineering represents only 50-65% of total cost. The rest \u2014 selling, scoping, legal, compliance, staffing, project management, knowledge transfer \u2014 is organizational friction. The firms announcing AI transformation by giving engineers Copilot are optimizing the minority of their cost structure while ignoring the majority.")


# ============================================================
# SLIDE 9 — Deal Size vs Profitability [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Minimize friction per unit of revenue and you can profitably serve $30K and $30M with the same margin")

# Two comparison cards
models = [
    ("Traditional model", "Overhead per deal is constant\n\u2192 Unprofitable below threshold\n\u2192 Optimized for deal SIZE", RED_ACCENT),
    ("Low-friction model", "Overhead scales with deal\n\u2192 Profitable at any size\n\u2192 Optimized for deal VELOCITY", GREEN_ACCENT),
]

for i, (title, desc, border_color) in enumerate(models):
    x = Inches(0.8) + i * Inches(5.8)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(5.3), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = border_color
    card.line.width = Pt(3)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.25)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.color.rgb = DARK_BG
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(16)
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = FONT

add_callout(slide, "Deal frequency is rising. Deal size is falling. The firm optimized for velocity has a structural advantage that compounds with every deal closed.", y=Inches(5.2))

add_speaker_notes(slide, "The implication is a fundamental reorientation. The economics of the old model \u2014 optimize delivery cost on a large, stable contract \u2014 don't transfer to a world where deal frequency is rising and deal size is falling. A firm that can move a $30K deal through its full pipeline at the same speed and margin as a $3M deal has a structural advantage. And that advantage compounds with every deal closed.")


# ============================================================
# SLIDE 10 — Incumbents Can't Buy Their Way Out [INVESTOR]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "$3B in AI spend doesn't fix a structural problem \u2014 Accenture's model requires the friction it's trying to eliminate")

responses = [
    ("Procure AI tools\nfor engineers", "Relocates bottleneck"),
    ("Rebrand revenue as\n\"AI consulting\"", "Optics, not structure"),
    ("Acquire AI startups", "Incompatible operating\nstructures"),
]

for i, (response, outcome) in enumerate(responses):
    x = Inches(0.8) + i * Inches(3.8)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.4), Inches(2.8))
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.25)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = response
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_BG
    run.font.name = FONT
    run.font.bold = True

    # Red X
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)
    run2 = p2.add_run()
    run2.text = "\u2717"
    run2.font.size = Pt(36)
    run2.font.color.rgb = RED_ACCENT
    run2.font.name = FONT

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(8)
    run3 = p3.add_run()
    run3.text = outcome
    run3.font.size = Pt(14)
    run3.font.color.rgb = MID_GRAY
    run3.font.name = FONT
    run3.font.italic = True

# Pull-quote
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "\"A container shipping company cannot pivot to same-day drone delivery by buying drones.\""
run.font.size = Pt(18)
run.font.color.rgb = ACCENT_TEAL
run.font.name = FONT
run.font.italic = True

add_footnote(slide, "Source: Accenture IR \u2014 $3B AI commitment, 31+ acquisitions/year, SG&A ~17% of revenue")

add_speaker_notes(slide, "Accenture has committed $3B to AI and acquires 31+ companies a year. Surely scale and capital can close this gap. It cannot. The problem is structural, not financial. An operation built to sell 10,000 junior developers at $30/hour is a volume business. The entire organizational infrastructure \u2014 campuses, training pipelines, bench management \u2014 exists to support labour arbitrage at scale. You cannot repurpose this into high-leverage AI-native delivery any more than a container shipping company can pivot to same-day drone delivery by buying drones.")


# ============================================================
# SLIDE 11 — Three-Layer Stack [ALL] (Section opener)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_title(slide, "Infrastructure commoditizes. Applications compress.\nThe context layer emerges as the value center.", color=WHITE)

# Three layers
layer_data = [
    ("LAYER 3 \u2014 CONTEXT LAYER", "Institutional knowledge directing AI agents", ACCENT_TEAL, Inches(2.4), Inches(1.4)),
    ("LAYER 2 \u2014 APPLICATIONS", "Dashboards, CRMs, project management\n$300B+ market cap lost. Growth halved 36% \u2192 17%.", MID_GRAY, Inches(4.0), Inches(1.2)),
    ("LAYER 1 \u2014 INFRASTRUCTURE", "APIs, databases, auth, monitoring\nCommoditizing. AI agents consume 100x more API calls.", RGBColor(0x44, 0x44, 0x66), Inches(5.4), Inches(1.0)),
]

for title, desc, color, y, h in layer_data:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), y, Inches(9.0), h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = WHITE if color == ACCENT_TEAL else LIGHT_GRAY
    run2.font.name = FONT

# Christensen's Law
txBox = slide.shapes.add_textbox(Inches(2.0), Inches(6.6), Inches(9.0), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Christensen's Law: when one layer commoditizes, adjacent layers de-commoditize."
run.font.size = Pt(14)
run.font.color.rgb = LIGHT_GRAY
run.font.name = FONT
run.font.italic = True

add_speaker_notes(slide, "Clayton Christensen's Law of Conservation of Attractive Profits: when one layer commoditizes, adjacent layers de-commoditize. Infrastructure is commoditizing \u2014 AI agents generate hundreds of API calls per minute versus a human generating dozens per hour. Applications are compressing \u2014 SaaS companies have lost over $300 billion in market cap, growth has halved, switching costs are dissolving as LLMs make data migration trivial. The context layer is what emerges. It's the institutional knowledge that directs AI agents \u2014 it doesn't live in a document. It emerges from thousands of workflows executed over time. This is where value concentrates.")


# ============================================================
# SLIDE 12 — AI Deployment Gap [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Firms deploy AI in products and for building products \u2014 almost none deploy AI to run the company")

modes = [
    ("AI in the product you sell", "\u2190 Most firms focus here", 0.85),
    ("AI for building the product", "\u2190 Second priority", 0.55),
    ("AI to run the company", "\u2190 Where structural economics change", 0.12),
]

for i, (mode, note, fill_pct) in enumerate(modes):
    y = Inches(2.2) + i * Inches(1.4)
    # Full bar background
    full_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(7.0), Inches(0.8))
    full_bar.fill.solid()
    full_bar.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    full_bar.line.fill.background()
    # Fill bar
    fill_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(7.0 * fill_pct), Inches(0.8))
    fill_bar.fill.solid()
    fill_bar.fill.fore_color.rgb = ACCENT_BLUE if i < 2 else ACCENT_TEAL
    fill_bar.line.fill.background()
    # Mode label
    txBox = slide.shapes.add_textbox(Inches(0.9), y + Inches(0.1), Inches(6.8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = mode
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE if fill_pct > 0.3 else DARK_TEXT
    run.font.name = FONT
    run.font.bold = True
    # Note
    txBox2 = slide.shapes.add_textbox(Inches(8.2), y + Inches(0.1), Inches(4.5), Inches(0.6))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    run = p.add_run()
    run.text = note
    run.font.size = Pt(14)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT

add_callout(slide, "If friction per unit of revenue is the metric, AI at one stage is not the solution. AI across the entire pipeline is.", y=Inches(6.0))

add_speaker_notes(slide, "There are three ways to deploy AI. In the product. For building the product. And to run the company. Most firms are focused on the first two. Very few are addressing the third. But the third is where the structural economics change. Sales must move faster. Pre-sales must scope in days, not weeks. Legal must review in hours. Compliance must clear in a single pass. Knowledge must transfer automatically, not through meetings. No IT services firm is structured to build this. Their teams are organized by function, not by flow. Their tooling is procured per function. There is no connective tissue.")


# ============================================================
# SLIDE 13 — realfast Operating Model [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Every function at realfast \u2014 pre-sales, delivery, operations, recruitment \u2014 runs on a single AI operating system")

functions = ["Pre-sales", "Scoping", "Legal", "Compliance", "Staffing", "Delivery", "QA", "PM"]
box_w = Inches(1.3)
box_h = Inches(0.6)
start_x = Inches(0.5)

# Traditional row
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(3.0), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "TRADITIONAL"
run.font.size = Pt(12)
run.font.color.rgb = MID_GRAY
run.font.name = FONT
run.font.bold = True

for i, fn in enumerate(functions):
    x = start_x + i * (box_w + Inches(0.12))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    box.line.color.rgb = MID_GRAY
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = fn
    run.font.size = Pt(10)
    run.font.color.rgb = DARK_TEXT
    run.font.name = FONT
    # Dashed arrows
    if i < len(functions) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w, Inches(2.35), Inches(0.12), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.95), Inches(11.0), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Siloed functions. Manual handoffs. Friction at every boundary."
run.font.size = Pt(12)
run.font.color.rgb = RED_ACCENT
run.font.name = FONT
run.font.italic = True

# realfast row
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(3.0), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "REALFAST"
run.font.size = Pt(12)
run.font.color.rgb = ACCENT_TEAL
run.font.name = FONT
run.font.bold = True

# Exocortex ribbon
ribbon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, Inches(4.6), Inches(11.5), Inches(0.5))
ribbon.fill.solid()
ribbon.fill.fore_color.rgb = ACCENT_TEAL
ribbon.line.fill.background()
tf = ribbon.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "E X O C O R T E X"
run.font.size = Pt(14)
run.font.color.rgb = WHITE
run.font.name = FONT
run.font.bold = True

for i, fn in enumerate(functions):
    x = start_x + i * (box_w + Inches(0.12))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.0), box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = ACCENT_TEAL
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = fn
    run.font.size = Pt(10)
    run.font.color.rgb = WHITE
    run.font.name = FONT

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(11.0), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Integrated pipeline. Context flows automatically. Friction structurally eliminated."
run.font.size = Pt(12)
run.font.color.rgb = GREEN_ACCENT
run.font.name = FONT
run.font.italic = True

add_speaker_notes(slide, "realfast is built on this thesis. Every function in the company runs on Exocortex \u2014 an AI operating system that spans the full pipeline. This is not a set of AI tools bolted onto existing processes. It is the operating infrastructure that connects everything. The connective tissue that traditional firms don't have.")


# ============================================================
# SLIDE 14 — Exocortex Functions [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Exocortex extracts institutional knowledge, maintains it as the business evolves, and applies it across every workflow")

# Flywheel nodes
nodes = [
    ("EXTRACT", "From people, contracts,\npast projects, domain expertise", ACCENT_BLUE),
    ("MAINTAIN", "Continuously re-engineer\nas we learn what works", ACCENT_TEAL),
    ("APPLY", "Across workflows \u2014\nwhoever needs it, whenever", RGBColor(0x66, 0x66, 0x99)),
]

for i, (title, desc, color) in enumerate(nodes):
    x = Inches(0.8) + i * Inches(4.0)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.2), Inches(3.5), Inches(2.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()

    tf = circle.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.color.rgb = WHITE
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = WHITE
    run2.font.name = FONT

    # Arrow between nodes
    if i < 2:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(3.5), Inches(2.9), Inches(0.5), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

# Center insight
txBox = slide.shapes.add_textbox(Inches(2.5), Inches(4.6), Inches(8.0), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Gets smarter by doing work, not by being trained."
run.font.size = Pt(22)
run.font.color.rgb = ACCENT_TEAL
run.font.name = FONT
run.font.bold = True

add_speaker_notes(slide, "Exocortex does three things. It extracts institutional knowledge from people, contracts, past projects, and domain expertise. It maintains that knowledge as the business evolves \u2014 continuously re-engineering processes as we learn what works. And it applies it across every workflow, for whoever needs it. Each engagement generates execution traces \u2014 granular records of what worked, what failed, what took longer than expected, what the client actually needed versus what the SOW described. These traces feed the next engagement. The system gets smarter by doing work, not by being trained.")


# ============================================================
# SLIDE 15 — Economic Mechanism [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Coordination overhead is 35-50% of delivery cost \u2014 Exocortex structurally eliminates it")

# Before column
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "BEFORE (Traditional)"
run.font.size = Pt(14)
run.font.color.rgb = RED_ACCENT
run.font.name = FONT
run.font.bold = True

roles_before = ["PM", "BA", "Architect", "Developer", "QA"]
for i, role in enumerate(roles_before):
    y = Inches(2.2) + i * Inches(0.65)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(2.0), Inches(0.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    box.line.color.rgb = MID_GRAY
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = role
    run.font.size = Pt(13)
    run.font.color.rgb = DARK_TEXT
    run.font.name = FONT
    if i < len(roles_before) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2.3), y + Inches(0.5), Inches(0.4), Inches(0.15))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_GRAY
        arrow.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(5.0), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Coordination overhead: 35-50%"
run.font.size = Pt(14)
run.font.color.rgb = RED_ACCENT
run.font.name = FONT
run.font.bold = True

# After column
txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "AFTER (realfast)"
run.font.size = Pt(14)
run.font.color.rgb = GREEN_ACCENT
run.font.name = FONT
run.font.bold = True

roles_after = ["Consultant", "Consultant"]
for i, role in enumerate(roles_after):
    y = Inches(2.6) + i * Inches(1.0)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), y, Inches(2.5), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = ACCENT_TEAL
    box.line.width = Pt(2)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = role
    run.font.size = Pt(14)
    run.font.color.rgb = WHITE
    run.font.name = FONT
    run.font.bold = True

# Exocortex connector
exo_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(4.2), Inches(4.0), Inches(0.5))
exo_bar.fill.solid()
exo_bar.fill.fore_color.rgb = ACCENT_TEAL
exo_bar.line.fill.background()
tf = exo_bar.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "EXOCORTEX"
run.font.size = Pt(13)
run.font.color.rgb = WHITE
run.font.name = FONT
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(6.8), Inches(5.0), Inches(5.5), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Coordination overhead: structurally eliminated"
run.font.size = Pt(14)
run.font.color.rgb = GREEN_ACCENT
run.font.name = FONT
run.font.bold = True

# Bottom hero number
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Not a 10% efficiency gain. A 2-3x margin improvement on transformed workflows."
run.font.size = Pt(20)
run.font.color.rgb = ACCENT_BLUE
run.font.name = FONT
run.font.bold = True

add_speaker_notes(slide, "The economic mechanism is specific. Coordination overhead in a typical project \u2014 project management, status alignment, knowledge transfer, cross-timezone handoffs, role-dependency wait times \u2014 represents 35-50% of total delivery cost. These aren't engineering costs. They're organizational costs. They show up in payroll budgets, not IT budgets. Exocortex doesn't optimize this overhead. It structurally eliminates it. Consultants who previously needed dedicated developers, PMs, and BAs can now deliver independently. The result is not a 10% gain. It's a 2-3x margin improvement.")


# ============================================================
# SLIDE 16 — Compounding [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Domain knowledge, workflow patterns, and playbooks compound in Exocortex \u2014 they don't walk out the door")

# Compounding curve (descending cost bars)
bar_heights = [Inches(2.0), Inches(1.5), Inches(1.1), Inches(0.8), Inches(0.55)]
labels = ["Project 1", "Project 2", "Project 3", "Project 4", "Project 5"]
base_y = Inches(5.0)
bar_w = Inches(1.5)

for i, (h, label) in enumerate(zip(bar_heights, labels)):
    x = Inches(1.0) + i * (bar_w + Inches(0.3))
    y = base_y - h
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    # Label below
    txBox = slide.shapes.add_textbox(x, base_y + Inches(0.1), bar_w, Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT

# Y-axis label
txBox = slide.shapes.add_textbox(Inches(0.2), Inches(2.8), Inches(0.8), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Cost /\nTime"
run.font.size = Pt(10)
run.font.color.rgb = MID_GRAY
run.font.name = FONT

# Flywheels on right
flywheels = [
    ("Data", "More deployments \u2192 better outcomes \u2192 more referrals"),
    ("Efficiency", "More usage \u2192 faster delivery \u2192 better margins"),
    ("Playbooks", "More patterns \u2192 higher capacity per person"),
]

for i, (title, desc) in enumerate(flywheels):
    y = Inches(2.0) + i * Inches(1.2)
    txBox = slide.shapes.add_textbox(Inches(8.5), y, Inches(4.5), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.color.rgb = ACCENT_TEAL
    run.font.name = FONT
    run.font.bold = True
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = FONT

add_callout(slide, "Traditional services firms scale linearly. We scale exponentially.", y=Inches(5.8))

add_speaker_notes(slide, "Every project completed makes the next one faster. In a traditional firm, knowledge dissipates when the team disbands. In our model, domain knowledge, workflow patterns, and client-specific context compound in Exocortex. The second project is faster than the first. The third is faster than the second. This is not true for traditional services firms. Their delivery quality is a function of whoever they happen to staff. Ours is a function of every engagement we've ever delivered.")


# ============================================================
# SLIDE 17 — Client Results [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Production results across industries \u2014 not demos, not pilots, not proofs of concept")

cases = [
    ("TASC", "3.5x", "Improvement in qualified leads", "799 \u2192 2,818 leads/month", "Featured Salesforce customer story"),
    ("Houzbay", "30%", "Conversion rate", "Nearly 1 in 3 users qualified as lead", ""),
    ("Flok", "5 days", "To production", "Complex Agentforce implementation", ""),
]

for i, (client, metric, label, detail, extra) in enumerate(cases):
    x = Inches(0.8) + i * Inches(3.9)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.5), Inches(3.0))
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = client
    run.font.size = Pt(16)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    run2 = p2.add_run()
    run2.text = metric
    run2.font.size = Pt(36)
    run2.font.color.rgb = ACCENT_BLUE
    run2.font.name = FONT
    run2.font.bold = True

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = label
    run3.font.size = Pt(14)
    run3.font.color.rgb = DARK_TEXT
    run3.font.name = FONT

    p4 = tf.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(6)
    run4 = p4.add_run()
    run4.text = detail
    run4.font.size = Pt(12)
    run4.font.color.rgb = MID_GRAY
    run4.font.name = FONT

    if extra:
        p5 = tf.add_paragraph()
        p5.alignment = PP_ALIGN.CENTER
        p5.space_before = Pt(4)
        run5 = p5.add_run()
        run5.text = extra
        run5.font.size = Pt(11)
        run5.font.color.rgb = ACCENT_TEAL
        run5.font.name = FONT
        run5.font.italic = True

# Quote
add_callout(slide, "\u201cIt's like hiring 3 quality SDRs to drive pipeline without adding headcount.\u201d \u2014 Sudheer Noohu, CTO, TASC", y=Inches(5.2))

add_speaker_notes(slide, "These are production results. TASC went from 799 leads in June to 2,818 in July \u2014 3.5x improvement while maintaining conversation quality. Their CTO said it's like hiring 3 quality SDRs without adding headcount. Houzbay's lead gen agent hit 30% conversion \u2014 nearly one in three users qualifying. Flok went from kickoff to production in 5 business days on a complex Agentforce implementation. These prove delivery speed and outcome quality. What they also prove \u2014 though it's not on the slide \u2014 is that each of these engagements fed the next. The compounding is already happening.")


# ============================================================
# SLIDE 18 — Team and Backing [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Founded by operators who've built at decacorn scale, backed by Peak XV and RTP Global")

founders = [
    ("Sidu Ponnappa", "CEO, Co-Founder", "[Background summary]"),
    ("Aakash Dharmadhikari", "CPTO, Co-Founder", "[Background summary]"),
    ("Steve Sule", "CFO, Co-Founder", "[Background summary]"),
]

for i, (name, role, bg) in enumerate(founders):
    x = Inches(0.8) + i * Inches(3.9)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.5), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.25)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = name
    run.font.size = Pt(20)
    run.font.color.rgb = DARK_BG
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    run2.text = role
    run2.font.size = Pt(14)
    run2.font.color.rgb = ACCENT_TEAL
    run2.font.name = FONT

    p3 = tf.add_paragraph()
    p3.space_before = Pt(8)
    run3 = p3.add_run()
    run3.text = bg
    run3.font.size = Pt(12)
    run3.font.color.rgb = MID_GRAY
    run3.font.name = FONT

# Investor logos bar
inv_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.2), Inches(8.0), Inches(0.7))
inv_bar.fill.solid()
inv_bar.fill.fore_color.rgb = DARK_BG
inv_bar.line.fill.background()
tf = inv_bar.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Peak XV Partners  |  RTP Global"
run.font.size = Pt(18)
run.font.color.rgb = WHITE
run.font.name = FONT
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Decades of software engineering, blitzscaling to decacorn scale, and billion-dollar M&A."
run.font.size = Pt(15)
run.font.color.rgb = MID_GRAY
run.font.name = FONT
run.font.italic = True

add_speaker_notes(slide, "[Personalize based on presenter. Key beats: relevant scale experience, domain credibility, why this team is uniquely positioned.]")


# ============================================================
# SLIDE 19 — Own Economics Placeholder [INVESTOR]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "[TO BE POPULATED] realfast unit economics vs. industry benchmarks")

placeholders = [
    "realfast RPE vs. industry benchmark ($50K)",
    "Margin profile by engagement type",
    "Revenue per FDE trend",
    "Cost-to-serve trend (declining)",
]

for i, text in enumerate(placeholders):
    y = Inches(2.0) + i * Inches(1.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(10.0), Inches(0.8))
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = MID_GRAY
    card.line.width = Pt(1)
    card.line.dash_style = 2  # Dashed

    tf = card.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"[DATA NEEDED] {text}"
    run.font.size = Pt(16)
    run.font.color.rgb = MID_GRAY
    run.font.name = FONT

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(6.2), Inches(10.0), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "PLACEHOLDER \u2014 Populate with actual data before presenting."
run.font.size = Pt(14)
run.font.color.rgb = RED_ACCENT
run.font.name = FONT
run.font.bold = True

add_speaker_notes(slide, "If we're making this argument about the industry, we need to show our own numbers. This is the most powerful proof point available. [Populate when data is ready.]")


# ============================================================
# SLIDE 20 — Pilot Proposal [ENTERPRISE]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "One workflow. One business metric. 3-6 weeks to production.")

steps = [
    ("1", "Start narrow", "One workflow, one metric"),
    ("2", "Prove the delta", "Time, cost, quality vs. traditional vendor"),
    ("3", "Expand", "Second project, third \u2014 each faster"),
    ("4", "End state", "Operational intelligence layer\ntuned to your business"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.6) + i * Inches(3.1)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), Inches(2.0), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_TEAL
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(22)
    run.font.color.rgb = WHITE
    run.font.name = FONT
    run.font.bold = True

    # Card below
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.0), Inches(2.8), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_BG
    run.font.name = FONT
    run.font.bold = True
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = MID_GRAY
    run2.font.name = FONT

    # Arrow between steps
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.8), Inches(3.7), Inches(0.3), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_TEAL
        arrow.line.fill.background()

# Risk box
risk_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.0))
risk_box.fill.solid()
risk_box.fill.fore_color.rgb = SUBTLE_BG
risk_box.line.color.rgb = GREEN_ACCENT
risk_box.line.width = Pt(2)
tf = risk_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Downside: "
run.font.size = Pt(14)
run.font.color.rgb = MID_GRAY
run.font.name = FONT
run2 = p.add_run()
run2.text = "3-6 weeks, minimal investment.   "
run2.font.size = Pt(14)
run2.font.color.rgb = DARK_TEXT
run2.font.name = FONT
run3 = p.add_run()
run3.text = "Upside: "
run3.font.size = Pt(14)
run3.font.color.rgb = MID_GRAY
run3.font.name = FONT
run4 = p.add_run()
run4.text = "compounding operational advantage."
run4.font.size = Pt(18)
run4.font.color.rgb = DARK_BG
run4.font.name = FONT
run4.font.bold = True

add_speaker_notes(slide, "We don't ask you to bet on a vision. We ask you to bet on one workflow. One business metric. 3-6 weeks. We prove the delta \u2014 time, cost, quality \u2014 against what you would have gotten from the firm you were about to hire. If we fail, you know quickly and cheaply. If we succeed, you have a decision to make about your next project. And every engagement compounds. The second project is faster than the first. The third is faster than the second. Over time, what you're building with us is not a series of disconnected implementations \u2014 it's an operational intelligence layer tuned to your business.")


# ============================================================
# SLIDE 21 — Competitive Urgency [ENTERPRISE]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "The firm that moves first compounds the advantage \u2014 the gap widens, not closes")

# Two diverging paths
# Competitor line (flat)
for i in range(6):
    x = Inches(1.0) + i * Inches(1.5)
    y = Inches(3.5)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = MID_GRAY
    dot.line.fill.background()
    if i < 5:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.2), y + Inches(0.08), Inches(1.3), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = MID_GRAY
        line.line.fill.background()

# Your line (ascending)
steps_y = [Inches(3.5), Inches(3.2), Inches(2.8), Inches(2.3), Inches(1.7), Inches(1.0)]
for i in range(6):
    x = Inches(1.0) + i * Inches(1.5)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, steps_y[i], Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_TEAL
    dot.line.fill.background()

# Labels
txBox = slide.shapes.add_textbox(Inches(9.5), Inches(3.3), Inches(3.0), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Competitor (Accenture)"
run.font.size = Pt(12)
run.font.color.rgb = MID_GRAY
run.font.name = FONT

txBox = slide.shapes.add_textbox(Inches(9.5), Inches(0.8), Inches(3.0), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "You (compounding)"
run.font.size = Pt(12)
run.font.color.rgb = ACCENT_TEAL
run.font.name = FONT
run.font.bold = True

# X-axis label
txBox = slide.shapes.add_textbox(Inches(4.0), Inches(4.2), Inches(5.0), Inches(0.3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Projects over time \u2192"
run.font.size = Pt(11)
run.font.color.rgb = MID_GRAY
run.font.name = FONT

add_callout(slide, "Nobody gets fired for hiring Accenture. But the firm that moves first compounds the advantage with every project.", y=Inches(5.0))

add_speaker_notes(slide, "Your competitors will hire Accenture. They'll get Accenture's cost structure, timeline, and coordination overhead. Nobody gets fired for that decision. But here's what that decision costs: every project they do with a traditional firm is a standalone engagement. No compounding. No operational intelligence accumulating. You, on the other hand, build compound advantage with every engagement. The gap doesn't stay constant. It widens.")


# ============================================================
# SLIDE 22 — Four Barriers to Replication [INVESTOR]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Competitors starting today are not behind on technology \u2014 they're behind on thousands of accumulated workflow traces")

barriers = [
    ("Extraction", "Context layer must be built through doing, not buying. Emerges from thousands of workflows executed."),
    ("Adoption", "Conviction drives adoption, not training. Deliver undeniable results first."),
    ("Compounding", "Every engagement generates execution traces encoding decision logic. Gap widens with every project."),
    ("Infrastructure mismatch", "Incumbents can't build this without dismantling what generates current revenue."),
]

for i, (title, desc) in enumerate(barriers):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(5.8)
    y = Inches(1.8) + row * Inches(2.2)

    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y, Inches(5.3), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(18)
    run.font.color.rgb = DARK_BG
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(14)
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = FONT

add_speaker_notes(slide, "This is not a technology moat \u2014 it's an operational moat. A competitor starting today is not just behind on code. They're behind on thousands of accumulated workflow traces that encode how workflows actually behave in practice. Which deals stall at legal review. Which client stakeholders need early alignment. Which compliance requirements are blocking versus procedural. This gap widens with every project we deliver. And the incumbents? They can't build this without dismantling the organizational infrastructure that generates their current revenue. The firms with the most to gain from this transition are the ones least able to execute it.")


# ============================================================
# SLIDE 23 — Client-Side Thesis [INVESTOR]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Quick wins through Exocortex compound into the client's own context layer")

progression = [
    ("1", "Deliver through Exocortex", "Client sees speed/cost delta"),
    ("2", "Work compounds", "Domain knowledge, workflows, context accumulate"),
    ("3", "Intelligence layer forms", "Client realizes operational leverage building"),
    ("4", "End state", "Client runs on same operating leverage as realfast"),
]

for i, (num, title, desc) in enumerate(progression):
    x = Inches(0.6) + i * Inches(3.1)
    # Step card — gets progressively wider
    w = Inches(2.5) + Inches(0.1 * i)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), w, Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = ACCENT_TEAL
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(24)
    run.font.color.rgb = ACCENT_TEAL
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = title
    run2.font.size = Pt(14)
    run2.font.color.rgb = DARK_BG
    run2.font.name = FONT
    run2.font.bold = True

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(6)
    run3 = p3.add_run()
    run3.text = desc
    run3.font.size = Pt(12)
    run3.font.color.rgb = MID_GRAY
    run3.font.name = FONT

    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w, Inches(2.8), Inches(0.3), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_TEAL
        arrow.line.fill.background()

add_callout(slide, "We don't sell Exocortex. We deliver through it. The platform thesis is an emergent outcome, not a sales pitch.", y=Inches(4.8))

add_speaker_notes(slide, "This is the second-order opportunity. Every engagement we deliver through Exocortex compounds context about the client's business. Over time, what they're building with us \u2014 without buying anything \u2014 is an operational intelligence layer tuned to their domain, their workflows, their institutional knowledge. The end state is not permanent dependency on realfast. It's their organization running on the same operating leverage we do. We don't sell Exocortex. We deliver through it. The platform thesis is an emergent outcome, not a sales pitch. This shifts our commercial model from subcontracting at better margins to transforming how services organizations operate.")


# ============================================================
# SLIDE 24 — Platform Stack [ALL]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "Exocortex is enterprise-grade, ISO 27001 and SOC 2 compliant, and operational at scale today")

components = [
    ("ExoCode", "AI-native\ndevelopment"),
    ("ExoWork", "Non-technical\nworkflows"),
    ("ExoManager", "Delivery\nmanagement"),
    ("ExoHelpdesk", "Managed\nservices"),
    ("ExoMetacog", "Agent\nanalytics"),
]

for i, (name, desc) in enumerate(components):
    x = Inches(0.5) + i * Inches(2.5)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.2), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.color.rgb = ACCENT_TEAL
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = name
    run.font.size = Pt(16)
    run.font.color.rgb = ACCENT_TEAL
    run.font.name = FONT
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(12)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.name = FONT

# Security badges
badge_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.4), Inches(6.0), Inches(0.7))
badge_bar.fill.solid()
badge_bar.fill.fore_color.rgb = SUBTLE_BG
badge_bar.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
badge_bar.line.width = Pt(1)
tf = badge_bar.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "ISO 27001   |   SOC 2"
run.font.size = Pt(18)
run.font.color.rgb = DARK_BG
run.font.name = FONT
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(2.0), Inches(5.5), Inches(9.0), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Not vaporware. Running our own operations at production scale."
run.font.size = Pt(18)
run.font.color.rgb = ACCENT_BLUE
run.font.name = FONT
run.font.bold = True

add_speaker_notes(slide, "This is not a prototype. Exocortex runs our operations today. Every engagement we've delivered \u2014 TASC, Houzbay, Flok \u2014 was delivered through this platform. It's ISO 27001 certified, SOC 2 compliant, and enterprise-grade. Five components working together: ExoCode for developers, ExoWork for business workflows, ExoManager for delivery oversight, ExoHelpdesk for managed services, ExoMetacog for analytics. All instrumented for full observability.")


# ============================================================
# SLIDE 25 — Investment Thesis [INVESTOR] (Dark)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_title(slide, "AI is restructuring IT services \u2014 realfast has built the platform to be the efficient winner", color=WHITE)

thesis = [
    ("What we are:", "A services company that delivers through a platform \u2014 every engagement compounds"),
    ("What compounds:", "Data, efficiency, playbooks, client context \u2014 all improve with every project"),
    ("What capital unlocks:", "Distribution (partner hires + acquisitions) to consolidate a fragmented market"),
]

for i, (label, desc) in enumerate(thesis):
    y = Inches(2.2) + i * Inches(1.3)
    txBox = slide.shapes.add_textbox(Inches(1.5), y, Inches(10.0), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{i+1}. {label} "
    run.font.size = Pt(18)
    run.font.color.rgb = ACCENT_TEAL
    run.font.name = FONT
    run.font.bold = True
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(18)
    run2.font.color.rgb = WHITE
    run2.font.name = FONT

# Exit paths
exits = ["Strategic acquisition", "PE roll-up", "Continue scaling"]
for i, exit_path in enumerate(exits):
    x = Inches(2.0) + i * Inches(3.3)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.0), Inches(2.8), Inches(0.5))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
    pill.line.color.rgb = ACCENT_TEAL
    pill.line.width = Pt(1)
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = exit_path
    run.font.size = Pt(13)
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = FONT

add_speaker_notes(slide, "Three things to remember. We are a services company that delivers through a platform \u2014 every engagement compounds. Capital unlocks distribution \u2014 partner hires with CXO rolodexes and SI acquisitions for client relationships. The exit math works three ways: strategic acquisition by a large SI that wants our operating model, PE roll-up of recurring managed services revenue, or continued scaling. The bet you're making is that AI will restructure IT services and we've built the platform to be the efficient winner in that restructuring.")


# ============================================================
# SLIDE 26 — Use of Funds [INVESTOR]
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

add_title(slide, "[TO BE POPULATED] Capital allocation and 12-18 month milestones")

# Allocation table
alloc_data = [
    ["Allocation", "Purpose"],
    ["Partner hires", "2-3 senior with CXO relationships"],
    ["Acquisition", "1-2 small SI acquisitions"],
    ["Platform", "Accelerate Exocortex roadmap"],
    ["Expansion", "Replicate Singapore model"],
]

table = slide.shapes.add_table(len(alloc_data), 2, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.2)).table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(3.5)

for row_idx, row in enumerate(alloc_data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        is_header = row_idx == 0
        set_cell_text(cell, val, font_size=13, bold=is_header or col_idx == 0, color=WHITE if is_header else DARK_TEXT)
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_WHITE

# Milestone table
milestone_data = [
    ["Milestone", "Target", "Timeframe"],
    ["Managed services ARR", "$[X]", "12 months"],
    ["Clients on platform", "[X]", "12 months"],
    ["First acquisition", "1", "18 months"],
]

table2 = slide.shapes.add_table(len(milestone_data), 3, Inches(6.8), Inches(1.6), Inches(5.8), Inches(1.8)).table
table2.columns[0].width = Inches(2.8)
table2.columns[1].width = Inches(1.5)
table2.columns[2].width = Inches(1.5)

for row_idx, row in enumerate(milestone_data):
    for col_idx, val in enumerate(row):
        cell = table2.cell(row_idx, col_idx)
        is_header = row_idx == 0
        set_cell_text(cell, val, font_size=13, bold=is_header or col_idx == 0, color=WHITE if is_header else DARK_TEXT)
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_WHITE

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "PLACEHOLDER \u2014 Finalize numbers before presenting."
run.font.size = Pt(14)
run.font.color.rgb = RED_ACCENT
run.font.name = FONT
run.font.bold = True

add_speaker_notes(slide, "[Populate with specific numbers. This slide converts conviction into commitment. Every number must be defensible.]")


# ============================================================
# SLIDE 27 — Enterprise Close [ENTERPRISE] (Dark)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, y=Inches(4.5), height=Inches(0.04))

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.0), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "One workflow. One metric. 3-6 weeks."
run.font.size = Pt(28)
run.font.color.rgb = WHITE
run.font.name = FONT
run.font.bold = True

txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10.0), Inches(0.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "If we don't deliver measurably better results, you've lost weeks \u2014 not months."
run.font.size = Pt(22)
run.font.color.rgb = LIGHT_GRAY
run.font.name = FONT

txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.0), Inches(0.8))
tf3 = txBox3.text_frame
p = tf3.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "That starts with one conversation."
run.font.size = Pt(30)
run.font.color.rgb = ACCENT_TEAL
run.font.name = FONT
run.font.bold = True

add_speaker_notes(slide, "One workflow. One metric. 3-6 weeks. We prove we can deliver the same outcome \u2014 production-grade, enterprise-quality \u2014 at a fundamentally different speed and cost. If we fail, you know in weeks with minimal investment. If we succeed, you have a measurable gap and a decision to make about your next project. There's no procurement cycle. No multi-year commitment. Just a focused test of whether a structurally different operating model delivers what the traditional firms can't.")


# ============================================================
# Save
# ============================================================
output_path = "/home/dev/workspace/steves-wiki/Comms/realfast_master_deck.pptx"
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
