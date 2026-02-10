from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUBTLE_BG = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
TABLE_ROW_ALT = RGBColor(0xF0, 0xF4, 0xF8)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BARRIER_BG = RGBColor(0xF8, 0xF0, 0xE0)
BARRIER_ACCENT = RGBColor(0xE8, 0x9D, 0x30)


def add_dark_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()


def add_light_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()


def add_accent_bar(slide, y=Inches(0), height=Inches(0.06)):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()


def set_cell_text(cell, text, font_size=11, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_footnote(slide, text, y=Inches(7.0)):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(8)
    run.font.color.rgb = MID_GRAY
    run.font.name = "Calibri"
    run.font.italic = True


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide)
add_accent_bar(slide, y=Inches(3.2), height=Inches(0.04))

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
for line in ["One AI execution partner.", "273 portfolio companies.", "Measurable impact in weeks."]:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    run = p.add_run()
    run.text = line
    run.font.size = Pt(36)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    run.font.bold = True
    p.space_after = Pt(8)

txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10), Inches(0.8))
tf2 = txBox2.text_frame
p = tf2.paragraphs[0]
run = p.add_run()
run.text = "realfast"
run.font.size = Pt(28)
run.font.color.rgb = ACCENT_TEAL
run.font.name = "Calibri"
run.font.bold = True
run2 = p.add_run()
run2.text = "  —  Portfolio-scale AI execution"
run2.font.size = Pt(20)
run2.font.color.rgb = LIGHT_GRAY
run2.font.name = "Calibri"

add_speaker_notes(slide, "None. Let the slide sit for 3-5 seconds. This is your opening image.")


# ============================================================
# SLIDE 2 — The Adoption Challenge
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

# Title
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The Adoption Challenge"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

# Subtitle
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Enterprise AI adoption is consistently lower than ambition.\nThe reasons are structural."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

# Bullets
bullets = [
    "Initiatives start too broad — no constrained scope",
    "AI is deployed as a tool, not embedded in how people work",
    "ROI is measured against assumptions, not operating metrics",
    "Hidden complexity erodes payback before value is visible",
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(3.0))
tf = txBox.text_frame
tf.word_wrap = True
for i, bullet in enumerate(bullets):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    run = p.add_run()
    run.text = bullet
    run.font.size = Pt(18)
    run.font.color.rgb = DARK_TEXT
    run.font.name = "Calibri"
    p.space_before = Pt(14)
    p.level = 0
    # Add bullet character
    run.text = "    " + bullet

add_footnote(slide, "40%+ of agentic AI projects will be canceled by 2027 — Gartner  |  Only 2% of organisations have deployed AI agents at full scale — Capgemini (April 2025)  |  Only 20% of EQT investment professionals generate the bulk of AI usage — EQT ThinQ, Q1 2025")

add_speaker_notes(slide, """EQT has declared AI adoption existential. You've built Motherbrain, partnered with OpenAI and Anthropic, hired 40 professionals in EQT Digital, and hosted an AI Summit for 75+ tech leaders across your portfolio. The ambition is clear.

But the industry data is equally clear. Over 40% of agentic AI projects will be canceled by 2027. Only 2% of organisations have deployed AI agents at full scale. And your own published data shows that even within EQT, 20% of investment professionals generate the bulk of AI usage.

The gap between ambition and adoption isn't about capability or intent. It's structural. Initiatives start too broad. AI gets bolted on rather than built into how people actually work. ROI gets measured against assumptions instead of real operating metrics. These are the same barriers across every enterprise we work with. And they're solvable — but they require a different execution model.""")


# ============================================================
# SLIDE 3 — The Portfolio Execution Challenge
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

# Title
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The Portfolio Execution Challenge"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

# Headline
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "EQT is all in on AI. The question is no longer whether — it's how to deploy it across 273 companies."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

# Three questions - as colored cards
questions = [
    ("Where", "Which workflow will produce the fastest, most measurable return?"),
    ("How", "Who executes it, and how fast can it reach production?"),
    ("What to measure", "What's the baseline, and what proves it worked?"),
]

card_width = Inches(3.5)
card_height = Inches(1.2)
start_x = Inches(0.8)
card_y = Inches(2.3)

for i, (label, desc) in enumerate(questions):
    x = start_x + i * (card_width + Inches(0.25))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(22)
    run.font.color.rgb = ACCENT_TEAL
    run.font.name = "Calibri"
    run.font.bold = True

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(13)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.name = "Calibri"
    p2.space_before = Pt(6)

# Five barriers
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Without answers, portfolio companies hit the same five barriers:"
run.font.size = Pt(15)
run.font.color.rgb = MID_GRAY
run.font.name = "Calibri"
run.font.bold = True

barriers = [
    ("No clear starting point", "budget exists, but no method to identify the highest-ROI workflow"),
    ("Scope too broad", '"AI transformation" instead of one workflow, one metric'),
    ("No baseline", "without a before measurement, there's no provable after"),
    ("Vendor fatigue", "6-12 month engagements that produce strategies, not systems"),
    ("Internal capacity", "IT teams built for maintenance, not AI-native delivery at speed"),
]

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, (title, desc) in enumerate(barriers):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{i+1}.  "
    run.font.size = Pt(14)
    run.font.color.rgb = ACCENT_BLUE
    run.font.name = "Calibri"
    run.font.bold = True
    run2 = p.add_run()
    run2.text = title
    run2.font.size = Pt(14)
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = "Calibri"
    run2.font.bold = True
    run3 = p.add_run()
    run3.text = f"  —  {desc}"
    run3.font.size = Pt(14)
    run3.font.color.rgb = MID_GRAY
    run3.font.name = "Calibri"
    p.space_before = Pt(8)

add_footnote(slide, "40%+ of agentic AI projects will be canceled by 2027 — Gartner  |  Only 2% of organisations have deployed AI agents at full scale — Capgemini (April 2025)")

add_speaker_notes(slide, """EQT has made its position clear — AI is existential, and the investment at the firm level reflects that. This isn't about whether AI matters. That question is settled.

The question now is: how does that conviction translate into production AI systems running inside the workflows of your portfolio companies? And when we look at what those companies face when they try, the same three questions come up every time.

Where do we start? Most companies have a mandate or budget, but no method to identify which specific workflow will produce the fastest, most measurable return. So the initiative starts too broad.

How do we execute? Even when they pick a target, the internal teams are built for maintenance, not for standing up an AI system in weeks. And external vendors are quoting 6 to 12 months and delivering strategy documents, not production systems.

What do we measure? This is the one that kills most projects. Nobody established a baseline before the work started, so when something gets built, there's no before-and-after. The CFO can't point to a number that moved. Without that, there's no case for expanding to the next workflow.

These barriers are structural. They repeat across every mid-market enterprise, regardless of sector. And they're solvable — but they require a specific execution model designed around answering those three questions. That's what we do.""")


# ============================================================
# SLIDE 4 — What We Do
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "What We Do"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "We are a horizontal AI execution partner. One proven model, across industries and use cases."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

# Three columns as cards
cols_data = [
    ("AUTOMATE", "Manual work creating backlogs and errors", "Extraction, validation, straight-through processing", "Cost reduction, throughput"),
    ("AUGMENT", "Decisions made without adequate support", "Scoring, recommendations, draft generation", "Revenue quality, decision speed"),
    ("CONNECT", "Systems and data that don't talk", "Cross-system matching, unified views", "Operational efficiency"),
]

card_w = Inches(3.6)
card_h = Inches(3.6)
start_x = Inches(0.8)
start_y = Inches(2.2)

for i, (title, problem, mechanism, lever) in enumerate(cols_data):
    x = start_x + i * (card_w + Inches(0.25))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.25)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.color.rgb = ACCENT_BLUE
    run.font.name = "Calibri"
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(14)
    run2 = p2.add_run()
    run2.text = problem
    run2.font.size = Pt(13)
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = "Calibri"

    p3 = tf.add_paragraph()
    p3.space_before = Pt(14)
    run3 = p3.add_run()
    run3.text = mechanism
    run3.font.size = Pt(12)
    run3.font.color.rgb = MID_GRAY
    run3.font.name = "Calibri"
    run3.font.italic = True

    p4 = tf.add_paragraph()
    p4.space_before = Pt(18)
    run4 = p4.add_run()
    run4.text = "EBITDA lever: "
    run4.font.size = Pt(12)
    run4.font.color.rgb = ACCENT_TEAL
    run4.font.name = "Calibri"
    run4.font.bold = True
    run5 = p4.add_run()
    run5.text = lever
    run5.font.size = Pt(12)
    run5.font.color.rgb = DARK_TEXT
    run5.font.name = "Calibri"

# Bottom line
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "These friction patterns repeat across every industry. The workflow details change. The underlying barriers are the same."
run.font.size = Pt(15)
run.font.color.rgb = MID_GRAY
run.font.name = "Calibri"
run.font.italic = True

add_speaker_notes(slide, """Before we apply AI to anything, we classify why work is slow. Every operational friction falls into one of three categories: work that should be automated, decisions that should be augmented, or systems that need to be connected.

What makes this powerful at portfolio scale is that these patterns repeat across every industry you invest in. A document processing backlog in healthcare is structurally identical to one in compliance services. A decision-latency problem in industrial scheduling is the same class of problem as one in financial underwriting. The workflow details change. The AI execution model doesn't. That's why we work horizontally.""")


# ============================================================
# SLIDE 5 — Why Horizontal Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Why Horizontal Works"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The same friction shows up everywhere. The industry wrapper changes — the core problem doesn't."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

# Matrix table
tbl_data = [
    ["Friction Pattern", "Healthcare", "Industrial Services", "Tech-Enabled Services", "Financial Services"],
    ["Execution backlog", "Patient intake", "Field service scheduling", "Client onboarding", "Claims processing"],
    ["Decision latency", "Clinical triage", "Maintenance scheduling", "Pricing and quoting", "Underwriting"],
    ["Cost-to-serve", "Staff scheduling", "Route optimisation", "Compliance checking", "Reconciliation"],
    ["Data disconnection", "Patient records", "Asset data across sites", "Client data across platforms", "Regulatory data"],
]

table = slide.shapes.add_table(len(tbl_data), len(tbl_data[0]), Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.2)).table

for col_idx in range(len(tbl_data[0])):
    table.columns[col_idx].width = Inches(2.3)

for row_idx, row in enumerate(tbl_data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        is_header = row_idx == 0
        is_first_col = col_idx == 0
        set_cell_text(
            cell, val,
            font_size=12 if not is_header else 13,
            bold=is_header or is_first_col,
            color=WHITE if is_header else (DARK_BG if is_first_col else DARK_TEXT)
        )
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_WHITE

# Bottom callout
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "One engagement teaches patterns that accelerate the next.\nBy the third company, we're faster, cheaper, and more accurate."
run.font.size = Pt(16)
run.font.color.rgb = ACCENT_BLUE
run.font.name = "Calibri"
run.font.bold = True

add_speaker_notes(slide, """This matrix is the key to why a horizontal model works better than vertical specialists at portfolio scale. Look at any row — the same structural problem appears in every sector. What we learn solving an execution backlog in healthcare directly accelerates how we solve one in industrial services.

By the third portfolio company engagement, we're meaningfully faster and cheaper — because we've already solved the structural problem elsewhere. That's the compounding advantage a horizontal partner delivers that a vertical specialist never can.""")


# ============================================================
# SLIDE 6 — Where We Start
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Where We Start"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Not every workflow is worth automating first. We start where the math works fastest."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

# Four criteria cards
criteria = [
    ("Direct business metric", "Something the CFO already tracks"),
    ("Repeatable workflow", "High frequency = fast proof"),
    ("Existing data", "No data engineering prerequisite"),
    ("Contained blast radius", "Failure stays local, success is attributable"),
]

card_w = Inches(2.6)
card_h = Inches(1.8)
start_x = Inches(0.8)

for i, (title, desc) in enumerate(criteria):
    x = start_x + i * (card_w + Inches(0.2))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.25)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_BG
    run.font.name = "Calibri"
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(13)
    run2.font.color.rgb = MID_GRAY
    run2.font.name = "Calibri"

# Entry points bar
bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.6), Inches(10.8), Inches(0.8))
bar.fill.solid()
bar.fill.fore_color.rgb = DARK_BG
bar.line.fill.background()
tf = bar.text_frame
tf.margin_left = Inches(0.3)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Entry points:   Execution backlogs   |   Decision latency   |   Cost-to-serve pressure   |   Revenue leakage"
run.font.size = Pt(16)
run.font.color.rgb = ACCENT_TEAL
run.font.name = "Calibri"
run.font.bold = True

add_speaker_notes(slide, """We're disciplined about where we start. Not every workflow is a good first target. We look for four things: a direct business metric the CFO already tracks, a repeatable workflow with enough volume to generate fast proof, existing data so we're not gated by a data engineering project, and a contained blast radius so that if it doesn't work, it doesn't cascade.

The entry points we're most confident in are execution backlogs — manual processing creating cost and delay; decision-to-execution latency — where the information exists but synthesis is manual; cost-to-serve pressure — where headcount keeps rising just to maintain throughput; and revenue leakage — manual quoting, slow responses, inconsistent pricing.

These aren't limitations. They're the doors that open everything else.""")


# ============================================================
# SLIDE 7 — One Playbook, Every Door
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "One Playbook, Every Door"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Same execution model. Regardless of industry, company, or use case."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

# Four steps
steps = [
    ("1", "Narrow scope", "One workflow, one metric, one team"),
    ("2", "Ship in weeks", "Production within 3-6 weeks"),
    ("3", "Measure against baseline", "Before/after on the metric that matters"),
    ("4", "Expand only from proof", "Next workflow earns its way in"),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(2.2) + i * Inches(1.1)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_TEAL
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(20)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    run.font.bold = True

    # Title + desc
    txBox = slide.shapes.add_textbox(Inches(2.0), y, Inches(9.0), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.color.rgb = DARK_BG
    run.font.name = "Calibri"
    run.font.bold = True
    run2 = p.add_run()
    run2.text = f"   {desc}"
    run2.font.size = Pt(16)
    run2.font.color.rgb = MID_GRAY
    run2.font.name = "Calibri"

# Bottom message
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "A repeatable model available to every portfolio company. Same playbook. Predictable timelines. Predictable costs. Measurable outcomes."
run.font.size = Pt(15)
run.font.color.rgb = ACCENT_BLUE
run.font.name = "Calibri"
run.font.bold = True

add_speaker_notes(slide, """This is how we stay horizontal without becoming vague. We don't need to understand your entire business. We need to understand one workflow deeply enough to make it measurably better. Then we move to the next.

For EQT, this means a model you can offer to portfolio companies across your sectors. The portfolio company gets a production AI system. You get a proven value creation lever with predictable costs, timelines, and measurable outcomes.""")


# ============================================================
# SLIDE 8 — What Portfolio Companies Receive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "What Portfolio Companies Receive"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Tangible outcomes. Not slide decks about AI."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

deliverables = [
    ("Production AI system", "Live, embedded in a real workflow, used by real people"),
    ("Measured business impact", "Before/after on the metric that justified the project"),
    ("Operational handover", "The company owns and maintains the system"),
    ("Expansion roadmap", "Where to go next, based on proven results"),
]

card_w = Inches(5.4)
card_h = Inches(1.0)

for i, (title, desc) in enumerate(deliverables):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * (card_w + Inches(0.3))
    y = Inches(2.2) + row * (card_h + Inches(0.3))

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = SUBTLE_BG
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_BG
    run.font.name = "Calibri"
    run.font.bold = True
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(13)
    run2.font.color.rgb = MID_GRAY
    run2.font.name = "Calibri"

# CFO line
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "What the CFO sees: "
run.font.size = Pt(16)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True
run2 = p.add_run()
run2.text = "a cost line that dropped, a throughput number that climbed, or a revenue metric that improved — with clear attribution."
run2.font.size = Pt(16)
run2.font.color.rgb = DARK_TEXT
run2.font.name = "Calibri"

add_speaker_notes(slide, """Every engagement produces four tangible deliverables. A production system that's live and being used. Measured before/after impact on the business metric. Full operational handover so the company owns it. And an expansion roadmap showing where to go next, based on what we've just proved.

What matters for EQT: this creates a repeatable proof point. When the next portfolio company asks 'does AI actually work for businesses like ours?' — you have a measured result from inside your own portfolio to point to.""")


# ============================================================
# SLIDE 9 — The Partnership Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, y=Inches(0))

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The Partnership Model"
run.font.size = Pt(32)
run.font.color.rgb = WHITE
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "How this works at portfolio scale."
run.font.size = Pt(18)
run.font.color.rgb = LIGHT_GRAY
run.font.name = "Calibri"

# EQT node at top
eqt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.9), Inches(3.0), Inches(0.7))
eqt_box.fill.solid()
eqt_box.fill.fore_color.rgb = ACCENT_TEAL
eqt_box.line.fill.background()
tf = eqt_box.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "EQT (Partnership)"
run.font.size = Pt(16)
run.font.color.rgb = WHITE
run.font.name = "Calibri"
run.font.bold = True

# Sector boxes
sectors = ["Healthcare", "Industrials", "Services", "..."]
for i, sector in enumerate(sectors):
    x = Inches(1.5) + i * Inches(2.8)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.2), Inches(2.2), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4E)
    box.line.color.rgb = ACCENT_BLUE
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sector
    run.font.size = Pt(14)
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = "Calibri"

    # realfast playbook box below
    if sector != "...":
        pb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.2), Inches(2.2), Inches(0.6))
        pb.fill.solid()
        pb.fill.fore_color.rgb = ACCENT_BLUE
        pb.line.fill.background()
        tf = pb.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "realfast Playbook"
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"
        run.font.bold = True

# Cross-portfolio intelligence bar
cpb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6))
cpb.fill.solid()
cpb.fill.fore_color.rgb = ACCENT_TEAL
cpb.line.fill.background()
tf = cpb.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Cross-Portfolio Intelligence"
run.font.size = Pt(16)
run.font.color.rgb = WHITE
run.font.name = "Calibri"
run.font.bold = True

# Flywheel steps
flywheel = [
    "1. Select together — identify companies where friction and readiness align",
    "2. We execute — same playbook, adapted to the workflow",
    "3. Results feed back — patterns, benchmarks, proven use cases",
    "4. Deploy to the next — each engagement accelerates the one after it",
]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(11.7), Inches(1.2))
tf = txBox.text_frame
tf.word_wrap = True
for i, step in enumerate(flywheel):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    run = p.add_run()
    run.text = step
    run.font.size = Pt(12)
    run.font.color.rgb = LIGHT_GRAY
    run.font.name = "Calibri"
    p.space_before = Pt(3)

add_speaker_notes(slide, """This is the model. We work together to identify portfolio companies where operational friction and readiness align. We execute using our standard playbook, adapted to the specific workflow. Results feed back — cross-portfolio pattern recognition, sector benchmarks, proven use cases. And then we deploy to the next company, where each engagement is faster and more informed because of the one before it.

What makes this compound over time: every engagement deepens the pattern recognition across the portfolio. This builds a cross-portfolio intelligence layer — benchmarks, proven use cases, sector-specific refinements — that gets more valuable with every deployment. That's a moat for both of us.""")


# ============================================================
# SLIDE 10 — The Economics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The Economics"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Phase model designed for PE risk appetite."
run.font.size = Pt(16)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

# Phase table
phase_data = [
    ["Phase", "What Happens", "Investment", "Timeline"],
    ["Prove", "1 company, 1 workflow", "$50-150K", "3-6 weeks"],
    ["Validate", "2-3 companies, cross-sector", "$150-400K", "6-12 weeks"],
    ["Scale", "Cross-portfolio deployment", "Framework pricing", "Ongoing"],
]

table = slide.shapes.add_table(len(phase_data), 4, Inches(0.8), Inches(1.6), Inches(7.0), Inches(1.4)).table
col_widths = [Inches(1.2), Inches(2.8), Inches(1.5), Inches(1.5)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

for row_idx, row in enumerate(phase_data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        is_header = row_idx == 0
        set_cell_text(cell, val, font_size=12, bold=is_header or col_idx == 0, color=WHITE if is_header else DARK_TEXT)
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_WHITE

# "We measure in operating metrics" callout
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "We measure in operating metrics, not aspirational percentages."
run.font.size = Pt(16)
run.font.color.rgb = ACCENT_BLUE
run.font.name = "Calibri"
run.font.bold = True

# Impact patterns table
impact_data = [
    ["What We Target", "Typical Before State", "Typical After State"],
    ["Manual recurring processes", "1-2 days of team time per cycle", "Under 1 hour per cycle"],
    ["Document processing throughput", "Hours per document, high error rate", "Minutes per document, consistent output"],
    ["Feature delivery on legacy systems", "Months-long backlogs, talent scarcity", "Weeks to production, existing codebase"],
    ["Revenue-critical deadlines", "Missed timelines, customer churn risk", "Deadlines met, accounts retained"],
]

table2 = slide.shapes.add_table(len(impact_data), 3, Inches(0.8), Inches(3.7), Inches(11.0), Inches(2.0)).table
impact_widths = [Inches(3.5), Inches(3.75), Inches(3.75)]
for i, w in enumerate(impact_widths):
    table2.columns[i].width = w

for row_idx, row in enumerate(impact_data):
    for col_idx, val in enumerate(row):
        cell = table2.cell(row_idx, col_idx)
        is_header = row_idx == 0
        set_cell_text(cell, val, font_size=12, bold=is_header or col_idx == 0, color=WHITE if is_header else DARK_TEXT)
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_WHITE

# Bottom callout
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "One workflow pays for itself. What follows is expansion — not a second pitch."
run.font.size = Pt(15)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True
p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Typical pattern: initial engagement surfaces 4-5 additional workflows.  Payback on first engagement: 2-4 months."
run2.font.size = Pt(13)
run2.font.color.rgb = MID_GRAY
run2.font.name = "Calibri"

add_footnote(slide, "Based on delivery patterns across engagements. Actual impact depends on workflow complexity, team size, and data readiness.")

add_speaker_notes(slide, """I'm not going to show you a slide that multiplies assumptions across 273 companies. You can do that math better than I can — you know your portfolio.

What I can show you is what happens at the individual company level when we execute. A finance team that was spending two days every month on manual reconciliation now spends less than an hour. A feature backlog on a legacy system that internal teams estimated at nine months gets compressed to weeks — because the alternative was losing their largest customer.

What's consistent across every engagement is the pattern afterward: one workflow proves the model, and that first engagement typically surfaces four or five additional workflows the company didn't realize were addressable. The compounding happens organically — we don't need to pitch a second engagement. The company comes to us.

The cost to prove this inside one portfolio company is $50 to $150K and 5-6 weeks. The payback on that first engagement is typically 2-4 months. After that, the question you'll be asking isn't whether this works — it's how many companies to deploy it across.""")


# ============================================================
# SLIDE 11 — Why realfast
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Why realfast"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Don't take our word for it. Test us."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

# Comparison table
comp_data = [
    ["", "realfast", "Perficient", "Big 4 / SI", "Boutique AI"],
    ["Built for", "Horizontal portfolio execution", "Broad digital transformation", "Enterprise advisory", "Deep tech R&D"],
    ["Speed", "3-6 weeks", "3-6 months", "6-12 months", "Varies"],
    ["Portfolio fit", "Same playbook, any industry", "Engagement by engagement", "Engagement by engagement", "Narrow"],
    ["Learning", "Compounds cross-portfolio", "Standalone", "Standalone", "No portfolio view"],
    ["Teams", "Senior engineers, small", "Mixed seniority, large", "Junior-heavy, large", "Research"],
    ["Pricing", "Outcome-anchored", "T&M, scope grows", "T&M, scope grows", "Project"],
]

table = slide.shapes.add_table(len(comp_data), 5, Inches(0.8), Inches(1.7), Inches(11.5), Inches(3.5)).table
comp_widths = [Inches(1.5), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)]
for i, w in enumerate(comp_widths):
    table.columns[i].width = w

for row_idx, row in enumerate(comp_data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        is_header = row_idx == 0
        is_realfast_col = col_idx == 1
        set_cell_text(
            cell, val, font_size=12,
            bold=is_header or col_idx == 0 or is_realfast_col,
            color=WHITE if is_header else (ACCENT_BLUE if is_realfast_col else DARK_TEXT)
        )
        if is_header:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEADER_BG
        elif is_realfast_col:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
        elif row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ROW_WHITE

# Bottom CTA
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.9))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Rather than debate it, let's prove it."
run.font.size = Pt(18)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True
p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Pick a portfolio company. Pick a workflow. Run us alongside your existing options. Same problem. Same timeline. Compare the output."
run2.font.size = Pt(15)
run2.font.color.rgb = MID_GRAY
run2.font.name = "Calibri"

add_speaker_notes(slide, """You've invested heavily in Perficient and that's the right call for large digital programs. We think we're the right call for fast, focused AI execution. But rather than debate the difference, let's demonstrate it.

Pick a portfolio company. Pick a workflow. Run us alongside your existing options — same problem, same timeline. Compare the output on speed, cost, adoption, and measured impact. We're confident enough to invite the comparison.""")


# ============================================================
# SLIDE 12 — How It Starts
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "How It Starts"
run.font.size = Pt(32)
run.font.color.rgb = DARK_BG
run.font.name = "Calibri"
run.font.bold = True

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "One decision. One company. One workflow."
run.font.size = Pt(18)
run.font.color.rgb = DARK_TEXT
run.font.name = "Calibri"

steps_data = [
    ("1", "Joint selection", "Identify a company based on friction and readiness"),
    ("2", "Scoping call", "1 hour. One workflow, one metric, success criteria."),
    ("3", "Discovery sprint", "2 weeks. Map workflow, establish baseline, design mechanism."),
    ("4", "Build + deploy", "3-4 weeks. Production system with real users."),
    ("5", "Measure + decide", "Hard numbers. Expand, or stop."),
]

for i, (num, title, desc) in enumerate(steps_data):
    y = Inches(2.0) + i * Inches(0.95)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(18)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    run.font.bold = True

    txBox = slide.shapes.add_textbox(Inches(2.0), y, Inches(9.5), Inches(0.55))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(18)
    run.font.color.rgb = DARK_BG
    run.font.name = "Calibri"
    run.font.bold = True
    run2 = p.add_run()
    run2.text = f"   {desc}"
    run2.font.size = Pt(15)
    run2.font.color.rgb = MID_GRAY
    run2.font.name = "Calibri"

# Timeline bar
bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3), Inches(10.8), Inches(0.7))
bar.fill.solid()
bar.fill.fore_color.rgb = DARK_BG
bar.line.fill.background()
tf = bar.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "5-6 weeks from handshake to measured impact."
run.font.size = Pt(20)
run.font.color.rgb = ACCENT_TEAL
run.font.name = "Calibri"
run.font.bold = True

add_speaker_notes(slide, """The engagement starts with one joint decision — which portfolio company, which workflow. A one-hour scoping call to define the metric and success criteria. Then a two-week discovery sprint to map the workflow, establish a baseline, and design the AI mechanism. Three to four weeks of build and deploy. And then we measure.

Total elapsed time: 5-6 weeks from handshake to measured impact. After Phase 1, the conversation changes from 'will this work?' to 'where do we deploy next?'""")


# ============================================================
# SLIDE 13 — Close
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, y=Inches(3.8), height=Inches(0.04))

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.6))
tf = txBox.text_frame
tf.word_wrap = True
lines = [
    "Enterprise AI adoption is held back by structural execution barriers —",
    "not by ambition or technology.",
    "",
    "A portfolio of 273+ companies represents both the challenge",
    "and the opportunity.",
]
for i, line in enumerate(lines):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    run = p.add_run()
    run.text = line
    run.font.size = Pt(22) if line else Pt(10)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10), Inches(0.8))
tf2 = txBox2.text_frame
p = tf2.paragraphs[0]
run = p.add_run()
run.text = "That starts with one conversation."
run.font.size = Pt(28)
run.font.color.rgb = ACCENT_TEAL
run.font.name = "Calibri"
run.font.bold = True

add_speaker_notes(slide, 'Let this slide breathe. Pause. Then: "We\'d like to explore what a Phase 1 engagement could look like with one of your portfolio companies. One conversation to see if the fit is there."')


# ============================================================
# Save
# ============================================================
output_path = "/home/dev/workspace/steves-wiki/EQT/realfast_x_EQT_v3.pptx"
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
